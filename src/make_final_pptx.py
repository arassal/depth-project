"""Build the final presentation deck for the depth-project paper.

Mirrors the paper structure (intro, related work, approach, experiments,
results, ablations, discussion, conclusion) and uses real numbers from
the outputs/ folder. Run from the repo root with the project venv:

    python src/make_final_pptx.py
"""

from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent.parent
OUT = ROOT / "outputs"
METRICS = OUT / "metrics"
FIGS = ROOT / "docs" / "assets" / "figures"
QUAL = FIGS / "qualitative"
SWEEPS = OUT / "plots" / "sweeps"
DOCS = ROOT / "docs"

OUT_PPTX = DOCS / "DepthProjectPresentation.pptx"
GITHUB = "https://github.com/arassal/depth-project"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text())


zero_nyu = load_json(METRICS / "zero_shot_nyu_summary.json")
distill_nyu = load_json(METRICS / "distill_test_summary.json")
zero_kitti = load_json(METRICS / "kitti_zero_shot_summary.json")
distill_kitti = load_json(METRICS / "kitti_distill_summary.json")


# Color palette: warm paper-like
BG = RGBColor(247, 240, 229)
INK = RGBColor(31, 26, 23)
MUTED = RGBColor(100, 90, 84)
ACCENT = RGBColor(168, 68, 37)
GOOD = RGBColor(56, 120, 70)
BAD = RGBColor(168, 50, 50)


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(1.0))
    tf = tx.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(30)
    r.font.bold = True
    r.font.color.rgb = INK
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(14)
        r2.font.color.rgb = MUTED


def add_bullets(slide, items, left, top, width, height, font_size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(item, tuple):
            text, level = item
        else:
            text, level = item, 0
        p.text = ("    " * level) + ("• " if level == 0 else "◦ ") + text
        p.font.size = Pt(font_size - 2 * level)
        p.font.color.rgb = color
        p.space_after = Pt(6)


def add_metric_box(slide, title, value, left, top, value_color=ACCENT, width=2.6, height=1.25):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 250, 241)
    shape.line.color.rgb = RGBColor(220, 205, 192)

    tf = shape.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    p1.text = title
    p1.font.size = Pt(12)
    p1.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.text = value
    p2.font.size = Pt(26)
    p2.font.bold = True
    p2.font.color.rgb = value_color


def add_caption(slide, text, left, top, width):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.4))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


def add_section_header(slide, label):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.0), Inches(12.2), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = label.upper()
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = ACCENT


def new_slide(section=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    if section:
        add_section_header(slide, section)
    return slide


# ---------------------------------------------------------------------------
# Slide 1: title (GitHub link + author names per the instructions)
# ---------------------------------------------------------------------------
slide = new_slide()
add_title(
    slide,
    "Distillation and Edge-Aware Loss for\nEfficient Monocular Depth Estimation",
    "Cal Poly Pomona  •  ECE 4990",
)

box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.7), Inches(2.5))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Alexander Assal  •  Parsa Ghasemi"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = INK
p2 = tf.add_paragraph()
p2.text = f"GitHub: {GITHUB}"
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT
p3 = tf.add_paragraph()
p3.text = "Student: Depth Anything Small  •  Teacher: Depth Anything Large  •  NYU-v2 + KITTI"
p3.font.size = Pt(14)
p3.font.color.rgb = MUTED


# ---------------------------------------------------------------------------
# Slide 2: outline
# ---------------------------------------------------------------------------
slide = new_slide("Outline")
add_title(slide, "Outline")
add_bullets(
    slide,
    [
        "Problem: monocular relative depth from a single RGB image",
        "Approach: distill Depth Anything Large → Small + multi-scale gradient loss",
        "Datasets: NYU-v2 (indoor) and KITTI eigen subset (outdoor)",
        "Three-axis hyperparameter sweep: distillation α, gradient β, learning rate",
        "Main finding: learning rate is the dominant hyperparameter",
        "Best stable config slightly beats zero-shot; aggressive lr collapses the student",
        "Live demo: local upload-and-predict web app",
    ],
    0.9, 1.5, 11.5, 5.0, font_size=20,
)


# ---------------------------------------------------------------------------
# Slide 3: problem
# ---------------------------------------------------------------------------
slide = new_slide("Background")
add_title(slide, "Monocular Depth is Ill-Posed")
add_bullets(
    slide,
    [
        "A single image is consistent with infinitely many world geometries",
        "Recovery is up to scale + shift unless extra information is supplied",
        "Modern systems train for relative depth, align scale + shift at eval",
        "Backbone: DPT head (Ranftl 2021) on DINOv2 ViT encoder (Oquab 2023)",
        "Foundation: Depth Anything (Yang 2024): scale + pseudo-labeled data",
    ],
    0.9, 1.5, 7.5, 4.5, font_size=20,
)
add_bullets(
    slide,
    [
        "Why this matters here:",
        "Depth Anything Small fits a 4 GB laptop GPU",
        "Depth Anything Large is the right teacher for distillation",
        "We want to push the Small student above its own pretrained baseline",
    ],
    8.6, 1.7, 4.3, 4.0, font_size=16, color=INK,
)


# ---------------------------------------------------------------------------
# Slide 4: related work
# ---------------------------------------------------------------------------
slide = new_slide("Related Work")
add_title(slide, "Related Work")
add_bullets(
    slide,
    [
        "Monocular depth (Eigen 2014, MiDaS 2022, ZoeDepth 2023)",
        "DPT dense prediction head (Ranftl 2021)",
        "Depth Anything v1/v2 (Yang 2024): 1.5M labeled + 62M pseudo-labeled",
        "DINOv2 backbone (Oquab 2023)",
        "Knowledge distillation (Hinton 2015)",
        "Distill-Any-Depth (He 2025): affine-invariant output-level distillation",
        "ZoeDepth multi-scale gradient L1 (Bhat 2023)",
    ],
    0.9, 1.5, 7.6, 4.8, font_size=18,
)
add_bullets(
    slide,
    [
        "What we lift:",
        "Affine-invariant output-level distillation → He 2025",
        "Multi-scale gradient L1 → Bhat 2023",
        "Pretrained checkpoints → Yang 2024 (HuggingFace)",
        "What we build new:",
        "Sweep harness, isolation, configs",
        "Three-axis controlled study (α, β, lr)",
    ],
    8.7, 1.5, 4.2, 5.0, font_size=14, color=INK,
)


# ---------------------------------------------------------------------------
# Slide 5: proposed approach overview
# ---------------------------------------------------------------------------
slide = new_slide("Proposed Approach")
add_title(slide, "Architecture", "Student Depth Anything Small  +  frozen Depth Anything Large teacher")
add_bullets(
    slide,
    [
        "Student: DINOv2 ViT-S/14 + DPT head, ~25M params, trainable",
        "Teacher: DINOv2 ViT-L/14 + DPT head, ~335M params, frozen",
        "Input: RGB 256×256, ImageNet mean/std normalization",
        "Output: relative depth, bilinearly upsampled to GT size",
        "All loss terms computed on aligned spatial size",
    ],
    0.9, 1.7, 11.5, 3.0, font_size=18,
)
add_bullets(
    slide,
    [
        "Why two models in parallel?",
        "The teacher provides a dense, low-noise per-pixel target on every training image",
        "Ground truth is sparse and noisy for far-field pixels; teacher complements it",
        "Frozen weights → deterministic supervision, no co-adaptation",
    ],
    0.9, 4.9, 11.5, 2.0, font_size=16, color=MUTED,
)


# ---------------------------------------------------------------------------
# Slide 6: loss function
# ---------------------------------------------------------------------------
slide = new_slide("Proposed Approach")
add_title(slide, "Loss Function")

box = slide.shapes.add_textbox(Inches(0.9), Inches(1.5), Inches(11.5), Inches(1.0))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "L_total = (1 – α) · L_base(student, GT)  +  α · L_distill(student, teacher)  +  β · L_grad(student, GT)"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(20)
p.font.bold = True
p.font.color.rgb = ACCENT

add_bullets(
    slide,
    [
        "L_base: affine-invariant L1 against NYU-v2 ground truth",
        "L_distill: affine-invariant L1 against frozen teacher prediction",
        "L_grad: multi-scale Sobel-style gradient L1 (scales 1, 2, 4)",
        "Robust normalization: per-sample median + MAD before L1",
        "Why MAD instead of std: stable under specular highlights, near-foreground outliers",
        "α = 0 → pure supervised fine-tune     α = 1 → pure distillation",
        "β controls boundary sharpness without affecting overall regression",
    ],
    0.9, 2.7, 11.5, 4.2, font_size=18,
)


# ---------------------------------------------------------------------------
# Slide 7: training procedure
# ---------------------------------------------------------------------------
slide = new_slide("Proposed Approach")
add_title(slide, "Training Procedure")
add_bullets(
    slide,
    [
        "Optimizer: AdamW, weight decay 1e-4, grad clip @ L2 norm 1.0",
        "Batch size: 1 (image_size=256 to fit 4 GB GPU)",
        "Epochs: 2 (small-data PoC regime)",
        "AMP mixed precision when CUDA available",
        "Teacher loaded once, frozen, never updated",
        "Validation AbsRel selects the best checkpoint",
    ],
    0.9, 1.5, 7.8, 4.5, font_size=18,
)
add_bullets(
    slide,
    [
        "Sweep harness:",
        "YAML grid → cartesian product",
        "Per-run isolated artifact folders",
        "Train + eval subprocesses",
        "Results aggregated to JSON",
    ],
    8.9, 1.7, 4.0, 4.0, font_size=15, color=INK,
)


# ---------------------------------------------------------------------------
# Slide 8: datasets
# ---------------------------------------------------------------------------
slide = new_slide("Experiments")
add_title(slide, "Datasets")
add_metric_box(slide, "NYU-v2 train", "450", 0.8, 1.2)
add_metric_box(slide, "NYU-v2 val", "60", 3.6, 1.2)
add_metric_box(slide, "NYU-v2 test", "59", 6.4, 1.2)
add_metric_box(slide, "KITTI test", "100", 9.2, 1.2)

add_bullets(
    slide,
    [
        "NYU-v2 PoC subset from sayakpaul/nyu_depth_v2 (Hugging Face)",
        "Float32 .npy depth in meters, range [0.001, 10.0] m",
        "KITTI eigen subset from exander/kitti-depth-gt (Hugging Face)",
        "Float32 .npy depth in meters, range [0.001, 80.0] m",
        "Both evaluated at 256×256 for internal comparability",
        "Test split sample IDs frozen across all runs",
    ],
    0.9, 3.0, 11.5, 4.0, font_size=18,
)


# ---------------------------------------------------------------------------
# Slide 9: metrics
# ---------------------------------------------------------------------------
slide = new_slide("Experiments")
add_title(slide, "Metrics")
add_bullets(
    slide,
    [
        "All metrics after per-image least-squares scale + shift alignment",
        "AbsRel  = mean(|pred – gt| / gt)     [↓ lower better]",
        "RMSE     = sqrt(mean((pred – gt)²))   [↓ lower better]",
        "RMSE_log = sqrt(mean((log pred – log gt)²))     [↓ lower better]",
        "log10    = mean(|log10 pred – log10 gt|)         [↓ lower better]",
        "δ_1, δ_2, δ_3 = fraction with max(p/t, t/p) < 1.25^n   [↑ higher better]",
        "Headline numbers reported: AbsRel and δ_1 (per Depth Anything paper)",
    ],
    0.9, 1.5, 11.5, 5.0, font_size=18,
)


# ---------------------------------------------------------------------------
# Slide 10: NYU results
# ---------------------------------------------------------------------------
slide = new_slide("Results")
add_title(slide, "NYU-v2 PoC Test Results", "59 images, image_size 256, after scale+shift alignment")

# Zero-shot
add_metric_box(slide, "Zero-shot AbsRel", f"{zero_nyu['abs_rel']:.4f}", 0.8, 1.3)
add_metric_box(slide, "Zero-shot δ_1", f"{zero_nyu['delta1']:.4f}", 3.55, 1.3)

# Default distill (collapse)
add_metric_box(slide, "Default AbsRel", f"{distill_nyu['abs_rel']:.4f}", 6.3, 1.3, value_color=BAD)
add_metric_box(slide, "Default δ_1", f"{distill_nyu['delta1']:.4f}", 9.05, 1.3, value_color=BAD)

# Best from sweep (improvement)
add_metric_box(slide, "Best AbsRel (ours)", "0.1505", 0.8, 2.85, value_color=GOOD)
add_metric_box(slide, "Best δ_1 (ours)", "0.8164", 3.55, 2.85, value_color=GOOD)
add_metric_box(slide, "AbsRel Δ vs ZS", "–0.0045", 6.3, 2.85, value_color=GOOD)
add_metric_box(slide, "δ_1 Δ vs ZS", "+0.0018", 9.05, 2.85, value_color=GOOD)

add_bullets(
    slide,
    [
        "Default config (α=0.5, β=0.1, lr=1e-5): student COLLAPSES — AbsRel doubles",
        "Best config (α=0.7, β=0.3, lr=1e-7): small but real improvement",
        "Improvement is consistent across AbsRel, RMSE, δ_1, δ_2, δ_3",
    ],
    0.9, 4.6, 11.5, 2.3, font_size=18,
)


# ---------------------------------------------------------------------------
# Slide 11: training collapse evidence
# ---------------------------------------------------------------------------
slide = new_slide("Results")
add_title(slide, "Training Collapse at Default Learning Rate", "Why we ran the sweep")

slide.shapes.add_picture(str(FIGS / "loss_vs_epoch.png"), Inches(0.5), Inches(1.5), width=Inches(4.0))
slide.shapes.add_picture(str(FIGS / "val_absrel_vs_epoch.png"), Inches(4.7), Inches(1.5), width=Inches(4.0))
slide.shapes.add_picture(str(FIGS / "val_delta1_vs_epoch.png"), Inches(8.9), Inches(1.5), width=Inches(4.0))
add_caption(slide, "Train loss", 0.5, 5.5, 4.0)
add_caption(slide, "Val AbsRel — returns to ~1.0 at epoch 2 (collapse)", 4.7, 5.5, 4.0)
add_caption(slide, "Val δ_1 — drops to 0.0 at epoch 2", 8.9, 5.5, 4.0)

add_bullets(
    slide,
    [
        "Signature: student output saturates to a near-constant value",
        "Triggered by (small batch=1) + (small training set 450 imgs) + (high lr 1e-5)",
        "Fix: drop the learning rate by 2 orders of magnitude → fully recovers",
    ],
    0.9, 6.0, 11.5, 1.4, font_size=15,
)


# ---------------------------------------------------------------------------
# Slide 12: alpha sweep
# ---------------------------------------------------------------------------
slide = new_slide("Ablation")
add_title(slide, "Distillation Weight α Sweep", "β=0.1, lr=1e-5 fixed")

slide.shapes.add_picture(str(SWEEPS / "abs_rel_vs_alpha.png"), Inches(0.5), Inches(1.4), width=Inches(6.2))
slide.shapes.add_picture(str(SWEEPS / "delta1_vs_alpha.png"), Inches(6.9), Inches(1.4), width=Inches(6.2))
add_caption(slide, "AbsRel vs α (↓ lower better)", 0.5, 5.8, 6.2)
add_caption(slide, "δ_1 vs α (↑ higher better)", 6.9, 5.8, 6.2)

add_bullets(
    slide,
    [
        "Sweep range: α ∈ {0.0, 0.3, 0.5, 0.7, 1.0}",
        "Best within sweep: α=0.7  →  AbsRel 0.202, δ_1 0.722",
        "Every point still WORSE than zero-shot (AbsRel 0.155) at this learning rate",
    ],
    0.9, 6.3, 11.5, 1.1, font_size=15,
)


# ---------------------------------------------------------------------------
# Slide 13: lr sweep
# ---------------------------------------------------------------------------
slide = new_slide("Ablation")
add_title(slide, "Learning Rate Sweep", "α=0.7, β=0.1 fixed   —   the dominant hyperparameter")

slide.shapes.add_picture(str(SWEEPS / "abs_rel_vs_learning_rate.png"), Inches(0.5), Inches(1.4), width=Inches(6.2))
slide.shapes.add_picture(str(SWEEPS / "delta1_vs_learning_rate.png"), Inches(6.9), Inches(1.4), width=Inches(6.2))
add_caption(slide, "AbsRel vs learning rate", 0.5, 5.8, 6.2)
add_caption(slide, "δ_1 vs learning rate", 6.9, 5.8, 6.2)

add_bullets(
    slide,
    [
        "lr=1e-7 → AbsRel 0.151    lr=1e-6 → 0.151    lr=5e-6 → 0.174    lr=1e-5 → 0.186",
        "Sharp transition between 1e-6 and 5e-6 — below = stable, above = drift",
        "Main paper finding: lr is the dominant lever, dominates α and β",
    ],
    0.9, 6.3, 11.5, 1.1, font_size=15,
)


# ---------------------------------------------------------------------------
# Slide 14: beta sweep
# ---------------------------------------------------------------------------
slide = new_slide("Ablation")
add_title(slide, "Gradient Weight β Sweep", "α=0.7, lr=1e-7 fixed   —   second-order effect")
slide.shapes.add_picture(str(SWEEPS / "abs_rel_vs_grad_weight.png"), Inches(1.5), Inches(1.4), width=Inches(10.0))
add_caption(slide, "AbsRel vs β (↓ lower better)", 1.5, 5.8, 10.0)
add_bullets(
    slide,
    [
        "β=0.0 → 0.152    β=0.1 → 0.152    β=0.3 → 0.150",
        "Direction is correct but the effect is small once lr is stabilized",
        "Gradient loss helps boundaries; quantitative impact is modest",
    ],
    0.9, 6.3, 11.5, 1.1, font_size=15,
)


# ---------------------------------------------------------------------------
# Slide 15: qualitative
# ---------------------------------------------------------------------------
slide = new_slide("Results")
add_title(slide, "Qualitative Examples", "NYU-v2 — RGB | GT | Pred | Error")
slide.shapes.add_picture(str(QUAL / "compare_000000.png"), Inches(0.4), Inches(1.4), width=Inches(6.3))
slide.shapes.add_picture(str(QUAL / "compare_000010.png"), Inches(6.8), Inches(1.4), width=Inches(6.3))
slide.shapes.add_picture(str(QUAL / "compare_000020.png"), Inches(0.4), Inches(3.9), width=Inches(6.3))
slide.shapes.add_picture(str(QUAL / "compare_000030.png"), Inches(6.8), Inches(3.9), width=Inches(6.3))
add_caption(slide, "Sample 000 — office scene", 0.4, 6.4, 6.3)
add_caption(slide, "Sample 010 — kitchen counter", 6.8, 6.4, 6.3)
add_bullets(
    slide,
    [
        "Scene structure preserved; error concentrates at depth discontinuities",
    ],
    0.9, 6.85, 11.5, 0.4, font_size=14, color=MUTED,
)


# ---------------------------------------------------------------------------
# Slide 16: KITTI
# ---------------------------------------------------------------------------
slide = new_slide("Cross-domain")
add_title(slide, "KITTI Cross-Domain Evaluation", "100 outdoor frames — trained only on NYU-v2 indoor")

add_metric_box(slide, "Zero-shot AbsRel", f"{zero_kitti['abs_rel']:.4f}", 0.8, 1.4)
add_metric_box(slide, "Distill AbsRel", f"{distill_kitti['abs_rel']:.4f}", 3.55, 1.4)
add_metric_box(slide, "Zero-shot δ_1", f"{zero_kitti['delta1']:.4f}", 6.3, 1.4)
add_metric_box(slide, "Distill δ_1", f"{distill_kitti['delta1']:.4f}", 9.05, 1.4)

slide.shapes.add_picture(str(FIGS / "kitti_input_0000000000.png"), Inches(1.6), Inches(3.2), width=Inches(4.5))
slide.shapes.add_picture(str(FIGS / "kitti_depth_0000000000.png"), Inches(7.2), Inches(3.2), width=Inches(4.5))
add_caption(slide, "KITTI RGB", 1.6, 6.4, 4.5)
add_caption(slide, "Predicted depth", 7.2, 6.4, 4.5)

add_bullets(
    slide,
    [
        "NYU-v2 improvement does NOT transfer to KITTI — essentially identical numbers",
        "Consequence of extreme-conservative recipe: encoder barely moves from pretrained init",
    ],
    0.9, 6.85, 11.5, 0.5, font_size=13, color=MUTED,
)


# ---------------------------------------------------------------------------
# Slide 17: discussion + limitations
# ---------------------------------------------------------------------------
slide = new_slide("Discussion")
add_title(slide, "Discussion")
add_bullets(
    slide,
    [
        "The dominant variable is the learning rate, not the loss balance",
        "Variance from α sweep: 0.09 AbsRel  vs  variance from lr sweep: 0.03 AbsRel",
        "But the lr range covers the difference between 'worse than ZS' and 'better than ZS'",
        "Improvement is small because lr=1e-7 keeps the encoder near initialization",
        "What would unlock more: larger effective batch, cross-context distillation (He 2025), more data",
    ],
    0.9, 1.5, 11.5, 3.5, font_size=18,
)
add_bullets(
    slide,
    [
        "Limitations:",
        "Evaluation at 256×256 inflates absolute KITTI numbers vs leaderboards",
        "Cross-domain test is a single OOD distribution",
        "Single GPU, single seed per grid point; no variance estimate",
    ],
    0.9, 5.2, 11.5, 2.0, font_size=15, color=MUTED,
)


# ---------------------------------------------------------------------------
# Slide 18: conclusion
# ---------------------------------------------------------------------------
slide = new_slide("Conclusion")
add_title(slide, "Conclusion")
add_bullets(
    slide,
    [
        "Built clean open-source pipeline: Depth Anything Small student + Large teacher",
        "Three losses: affine-invariant base, affine-invariant distill, multi-scale gradient",
        "Three-axis sweep: α (distill weight), β (grad weight), learning rate",
        "Empirically tuned config slightly beats pretrained zero-shot on NYU-v2",
        "No cross-domain transfer to KITTI at our conservative learning rate",
        "Headline finding: lr-sensitivity, not loss balance, dominates small-data distillation",
    ],
    0.9, 1.5, 11.5, 4.5, font_size=18,
)
add_bullets(
    slide,
    [
        f"GitHub: {GITHUB}",
        "Demo: ./run_upload_demo.sh  →  http://127.0.0.1:8000",
    ],
    0.9, 6.0, 11.5, 1.2, font_size=16, color=ACCENT,
)


# ---------------------------------------------------------------------------
# Slide 19: contributions + thanks
# ---------------------------------------------------------------------------
slide = new_slide("Acknowledgements")
add_title(slide, "Contributions  &  Thanks")
add_bullets(
    slide,
    [
        "Alexander Assal: pipeline implementation, distillation hook, gradient loss, sweep harness, web demo, analysis",
        "Parsa Ghasemi: literature survey, dataset preparation, qualitative analysis, figures, manuscript drafting",
        "Both authors: manuscript & presentation",
    ],
    0.9, 1.5, 11.5, 2.3, font_size=18,
)
add_bullets(
    slide,
    [
        "References we depend on:",
        "Depth Anything — Yang et al., CVPR 2024",
        "Distill-Any-Depth — He et al., arXiv 2025",
        "ZoeDepth — Bhat et al., arXiv 2023",
        "DPT — Ranftl et al., ICCV 2021",
        "DINOv2 — Oquab et al., arXiv 2023",
        "NYU-v2 — Silberman et al., ECCV 2012   •   KITTI — Geiger et al., IJRR 2013",
    ],
    0.9, 4.1, 11.5, 3.2, font_size=15, color=MUTED,
)

# Q&A slide
slide = new_slide()
add_title(slide, "Questions?", "Alexander Assal  •  Parsa Ghasemi")
box = slide.shapes.add_textbox(Inches(0.9), Inches(3.5), Inches(11.5), Inches(2.0))
tf = box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.text = GITHUB
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = ACCENT


OUT_PPTX.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT_PPTX)
print(f"wrote {OUT_PPTX}")
