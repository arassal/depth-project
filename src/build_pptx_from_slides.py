"""Build DepthProjectPresentation.pptx by embedding each Beamer slide as a
full-bleed image. Run after rebuilding slides.pdf:

    pdftoppm -r 200 paper/slides.pdf /tmp/slidespng/s -png
    python src/build_pptx_from_slides.py

Output: docs/DepthProjectPresentation.pptx (16:9, 21 slides)
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

PNG_DIR = Path("/tmp/slidespng")
OUT = Path(__file__).resolve().parent.parent / "docs" / "DepthProjectPresentation.pptx"

prs = Presentation()
# 16:9 (matches the Beamer aspectratio=169)
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

pngs = sorted(PNG_DIR.glob("s-*.png"))
assert pngs, f"no PNGs found in {PNG_DIR}"

blank = prs.slide_layouts[6]
for png in pngs:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0, width=prs.slide_width, height=prs.slide_height)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(f"wrote {OUT} ({len(pngs)} slides)")
