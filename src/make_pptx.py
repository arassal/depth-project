from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path("/home/alexander/depth-project")
OUT = ROOT / "outputs"
DESKTOP = Path("/home/alexander/Desktop/DepthProjectDemo")

METRICS = OUT / "metrics"
PLOTS = OUT / "plots"
PREDS = OUT / "predictions"
FIGS = ROOT / "docs" / "assets" / "figures"


def load_json(path: Path) -> dict:
    return json.loads(path.read_text())


zero = load_json(METRICS / "poc_zero_shot_summary.json")
supervised = load_json(METRICS / "poc_finetuned_summary.json")
student = load_json(METRICS / "teacher_student_poc_test_summary.json")
refined = load_json(METRICS / "teacher_student_refined_test_summary.json")
comparison = load_json(METRICS / "refined_architecture_comparison.json")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(247, 240, 229)
INK = RGBColor(31, 26, 23)
MUTED = RGBColor(100, 90, 84)
ACCENT = RGBColor(168, 68, 37)
OK = RGBColor(45, 106, 79)


def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None):
    tx = slide.shapes.add_textbox(Inches(0.55), Inches(0.3), Inches(12.2), Inches(1.0))
    p = tx.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(27)
    r.font.bold = True
    r.font.color.rgb = INK
    if subtitle:
        p2 = tx.text_frame.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(13)
        r2.font.color.rgb = MUTED


def add_bullets(slide, items: list[str], left: float, top: float, width: float, height: float, font_size: int = 20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = INK
        p.space_after = Pt(8)


def add_metric_box(slide, title: str, value: str, left: float, top: float, width: float = 2.25):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(1.25))
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 250, 241)
    shape.line.color.rgb = RGBColor(220, 205, 192)
    tf = shape.text_frame
    p1 = tf.paragraphs[0]
    p1.text = title
    p1.font.size = Pt(13)
    p1.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    p2.text = value
    p2.font.size = Pt(23)
    p2.font.bold = True
    p2.font.color.rgb = ACCENT


def add_caption(slide, text: str, left: float, top: float, width: float):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(0.5))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(11)
    p.font.color.rgb = MUTED


# Slide 1
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Depth Anything Teacher-Student Proof of Concept", "Alexander Assal • GitHub: https://github.com/arassal/depth-project")
add_bullets(
    slide,
    [
        "Goal: match the original paper story more closely and add a real student-side architecture change",
        "Teacher: pretrained Depth Anything Small",
        "Student: trained on labeled NYU-v2 subset plus teacher-pseudo-labeled indoor images",
        "Architecture change: RGB-guided residual refinement head",
        "Main result: refined student improved over earlier student runs but did not beat the teacher",
    ],
    0.8,
    1.4,
    11.6,
    3.0,
)

# Slide 2
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Student Architecture Change")
add_bullets(
    slide,
    [
        "Pretrained model: LiheYoung/depth-anything-small-hf",
        "Architecture family: Depth Anything Small / DPT-style monocular depth estimation",
        "New student head input: RGB plus coarse predicted depth",
        "New student head: 3 convolution layers with GELU",
        "Final output: coarse depth plus scaled residual correction",
    ],
    0.8,
    1.3,
    6.0,
    4.0,
)
add_bullets(
    slide,
    [
        "The teacher stays the original Depth Anything Small model",
        "The student gets a lightweight refinement module instead of a full backbone rewrite",
    ],
    7.0,
    2.0,
    5.1,
    2.0,
    font_size=18,
)

# Slide 3
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Data Growth")
add_metric_box(slide, "Labeled seed", "450", 0.8, 1.2)
add_metric_box(slide, "Pseudo labels", "2000", 3.25, 1.2)
add_metric_box(slide, "Student total", "2450", 5.7, 1.2)
add_metric_box(slide, "Held-out test", "59", 8.15, 1.2)
slide.shapes.add_picture(str(FIGS / "data_growth.png"), Inches(0.9), Inches(3.0), width=Inches(5.5))
add_bullets(
    slide,
    [
        "The student now trains on materially more data than the original labeled subset",
        "This is the closest local version of the paper's teacher-student data story",
    ],
    7.0,
    3.3,
    5.2,
    2.4,
    font_size=18,
)

# Slide 4
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Teacher-Student Workflow")
add_bullets(
    slide,
    [
        "1. Export 2000 indoor RGB images into an unlabeled pool",
        "2. Run the pretrained teacher on that pool",
        "3. Save teacher pseudo-depth maps",
        "4. Train the student on 450 labeled plus 2000 pseudo-labeled images",
        "5. Evaluate teacher, supervised-only student, and teacher-student student on the same test split",
    ],
    0.9,
    1.4,
    11.4,
    4.2,
)

# Slide 5
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Four-Run Results")
add_metric_box(slide, "Teacher AbsRel", f"{zero['abs_rel']:.4f}", 0.8, 1.2)
add_metric_box(slide, "Teacher delta1", f"{zero['delta1']:.4f}", 3.25, 1.2)
add_metric_box(slide, "Supervised AbsRel", f"{supervised['abs_rel']:.4f}", 5.7, 1.2)
add_metric_box(slide, "Refined AbsRel", f"{refined['abs_rel']:.4f}", 8.15, 1.2)
slide.shapes.add_picture(str(FIGS / "refined_run_comparison.png"), Inches(0.85), Inches(2.8), width=Inches(5.7))
add_bullets(
    slide,
    [
        f"Teacher-student vs supervised: AbsRel {comparison['teacher_student_student']['abs_rel'] - comparison['supervised_only_student']['abs_rel']:+.4f}, delta1 {comparison['teacher_student_student']['delta1'] - comparison['supervised_only_student']['delta1']:+.4f}",
        f"Refined vs teacher-student: AbsRel {comparison['refined_vs_teacher_student']['delta_abs_rel']:+.4f}, delta1 {comparison['refined_vs_teacher_student']['delta_delta1']:+.4f}",
        f"Refined vs teacher: AbsRel {comparison['refined_vs_zero_shot']['delta_abs_rel']:+.4f}, delta1 {comparison['refined_vs_zero_shot']['delta_delta1']:+.4f}",
    ],
    6.8,
    3.2,
    5.5,
    2.6,
    font_size=18,
)

# Slide 6
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Teacher Pseudo Labels")
slide.shapes.add_picture(str(FIGS / "teacher_pseudo_00.png"), Inches(0.8), Inches(1.4), width=Inches(5.7))
slide.shapes.add_picture(str(FIGS / "teacher_pseudo_01.png"), Inches(6.85), Inches(1.4), width=Inches(5.7))
add_caption(slide, "Teacher pseudo-depth example 1", 0.8, 6.7, 5.7)
add_caption(slide, "Teacher pseudo-depth example 2", 6.85, 6.7, 5.7)

# Slide 7
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Qualitative Sample 000000")
slide.shapes.add_picture(str(FIGS / "nyu_zero_shot_000000.png"), Inches(0.2), Inches(1.35), width=Inches(3.1))
slide.shapes.add_picture(str(FIGS / "nyu_finetuned_000000.png"), Inches(3.45), Inches(1.35), width=Inches(3.1))
slide.shapes.add_picture(str(FIGS / "teacher_student_000000.png"), Inches(6.7), Inches(1.35), width=Inches(3.1))
slide.shapes.add_picture(str(FIGS / "teacher_student_refined_000000.png"), Inches(9.95), Inches(1.35), width=Inches(3.1))
add_caption(slide, "Zero-shot teacher", 0.2, 6.55, 3.1)
add_caption(slide, "Supervised-only student", 3.45, 6.55, 3.1)
add_caption(slide, "Teacher-student student", 6.7, 6.55, 3.1)
add_caption(slide, "Refined student", 9.95, 6.55, 3.1)

# Slide 8
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Refined Training Curves")
slide.shapes.add_picture(str(FIGS / "refined_loss_vs_epoch.png"), Inches(0.7), Inches(1.3), width=Inches(4.0))
slide.shapes.add_picture(str(FIGS / "refined_val_absrel_vs_epoch.png"), Inches(4.75), Inches(1.3), width=Inches(4.0))
slide.shapes.add_picture(str(FIGS / "refined_val_delta1_vs_epoch.png"), Inches(8.8), Inches(1.3), width=Inches(4.0))
add_caption(slide, "Loss vs epoch", 0.7, 5.4, 4.0)
add_caption(slide, "Validation AbsRel", 4.75, 5.4, 4.0)
add_caption(slide, "Validation delta1", 8.8, 5.4, 4.0)

# Slide 9
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "How To Interpret The Outcome")
add_bullets(
    slide,
    [
        "The student is better with pseudo-labeled expansion than with the tiny supervised-only run",
        "The residual refinement head improves the student again over the plain teacher-student run",
        "The zero-shot teacher still generalizes best on the held-out test subset",
        "This means the local teacher-student pipeline and the architecture change are both valid, but the student still needs either more data or safer tuning to surpass the teacher",
    ],
    0.9,
    1.5,
    11.2,
    3.0,
)
add_metric_box(slide, "Teacher", "best", 1.1, 5.0, width=2.6)
add_metric_box(slide, "Supervised-only", "worst", 4.05, 5.0, width=2.8)
shape = slide.shapes.add_shape(1, Inches(7.2), Inches(5.0), Inches(4.1), Inches(1.25))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 250, 241)
shape.line.color.rgb = RGBColor(220, 205, 192)
tf = shape.text_frame
p1 = tf.paragraphs[0]
p1.text = "Teacher-student"
p1.font.size = Pt(13)
p1.font.color.rgb = MUTED
p2 = tf.add_paragraph()
p2.text = "middle, improved"
p2.font.size = Pt(22)
p2.font.bold = True
p2.font.color.rgb = OK

# Slide 10
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Second Dataset and Live Demo")
add_bullets(
    slide,
    [
        "KITTI Tiny is included as an outdoor zero-shot demo",
        "The uploadable local web app runs the pretrained teacher on arbitrary images",
        "That is the best live demo mode because the teacher remains the strongest model",
    ],
    0.8,
    1.2,
    6.0,
    2.6,
)
slide.shapes.add_picture(str(FIGS / "kitti_input_0000000000.png"), Inches(0.9), Inches(4.0), width=Inches(5.2))
slide.shapes.add_picture(str(FIGS / "kitti_depth_0000000000.png"), Inches(7.1), Inches(4.0), width=Inches(5.2))
add_caption(slide, "KITTI RGB sample", 0.9, 6.7, 5.2)
add_caption(slide, "KITTI predicted depth", 7.1, 6.7, 5.2)

# Slide 11
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(slide)
add_title(slide, "Final Takeaway")
add_bullets(
    slide,
    [
        "The project now matches the original paper story much more closely",
        "Training data expanded through pseudo labels and the student gained a refinement head",
        "The refined student improved over the earlier student baselines",
        "The pretrained teacher still remained strongest on the held-out test split",
    ],
    0.8,
    1.7,
    11.6,
    3.0,
)


out_path = DESKTOP / "DepthProjectPresentation.pptx"
out_path.parent.mkdir(parents=True, exist_ok=True)
prs.save(out_path)
print(out_path)
